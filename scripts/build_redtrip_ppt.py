#!/usr/bin/env python3
"""Build RedTrip contest PPT — 总分总结构，海派宣纸美学。"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "pptx"
OUT = ROOT / "docs" / "pptx" / "RedTrip-产品介绍.pptx"
ART = Path("/opt/cursor/artifacts/pptx/RedTrip-产品介绍.pptx")

# 海派 6 色
INK = RGBColor(0x33, 0x33, 0x3A)
INK_DEEP = RGBColor(0x1A, 0x1C, 0x21)
OCHRE = RGBColor(0xB9, 0x82, 0x4F)
SLATE = RGBColor(0x7C, 0x8A, 0x8D)
RICE = RGBColor(0xED, 0xE4, 0xD3)
VERMILION = RGBColor(0xA8, 0x32, 0x2A)
XUAN = RGBColor(0xF2, 0xEB, 0xDD)
PAPER = RGBColor(0xFB, 0xF8, 0xF1)
MUTED = RGBColor(0x6A, 0x61, 0x54)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def _set_run_font(run, *, size=18, bold=False, color=INK_DEEP, name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", name)


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sh, color)
    return sh


def _round(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _fill(sh, color)
    return sh


def _text(slide, l, t, w, h, text, *, size=18, bold=False, color=INK_DEEP, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, bold=bold, color=color)
    return box


def _para(tf, text, *, size=16, bold=False, color=INK_DEEP, align=PP_ALIGN.LEFT, space_before=0, space_after=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run_font(run, size=size, bold=bold, color=color)
    return p


def _bg(slide):
    _rect(slide, 0, 0, W, H, PAPER)
    # top wash strip
    _rect(slide, 0, 0, W, Inches(0.08), VERMILION)
    # bottom rice line
    _rect(slide, 0, H - Inches(0.35), W, Inches(0.35), XUAN)


def _footer(slide, page: str, section: str = "红鸢 RedTrip"):
    _text(slide, Inches(0.5), H - Inches(0.3), Inches(8), Inches(0.25), f"{section}  ·  城市记忆策展人", size=11, color=SLATE)
    _text(slide, W - Inches(1.4), H - Inches(0.3), Inches(1), Inches(0.25), page, size=11, color=SLATE, align=PP_ALIGN.RIGHT)


def _icon(slide, name: str, l, t, size=Inches(0.7)):
    path = ASSETS / name
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=size, height=size)


def _card(slide, l, t, w, h):
    sh = _round(slide, l, t, w, h, WHITE)
    sh.line.color.rgb = RICE
    sh.line.width = Pt(1)
    return sh


def _section_title(slide, eyebrow: str, title: str, subtitle: str = ""):
    _text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.3), eyebrow, size=14, color=VERMILION, bold=True)
    _text(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.55), title, size=32, color=INK_DEEP, bold=True)
    if subtitle:
        _text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4), subtitle, size=15, color=MUTED)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    page = 0

    def pnum():
        nonlocal page
        page += 1
        return str(page).zfill(2)

    # ── 1 Cover ──
    s = new_slide(prs)
    _rect(s, 0, 0, W, H, XUAN)
    _rect(s, 0, 0, W, Inches(0.12), VERMILION)
    _rect(s, 0, H - Inches(1.1), W, Inches(1.1), INK_DEEP)
    if (ASSETS / "redtrip-logo.jpg").exists():
        s.shapes.add_picture(str(ASSETS / "redtrip-logo.jpg"), Inches(0.7), Inches(1.3), width=Inches(5.2))
    _icon(s, "ppt-icon-kite.png", Inches(7.2), Inches(1.4), Inches(1.3))
    _icon(s, "ppt-icon-seal.png", Inches(8.7), Inches(1.55), Inches(0.9))
    _text(s, Inches(7.1), Inches(2.9), Inches(5.5), Inches(0.7), "红鸢 RedTrip", size=40, bold=True, color=VERMILION)
    _text(s, Inches(7.1), Inches(3.55), Inches(5.5), Inches(0.4), "城市记忆策展人", size=22, color=INK_DEEP)
    _text(s, Inches(7.1), Inches(4.15), Inches(5.5), Inches(0.8),
          "牵引城市深处的记忆之线\n化作一场可以行走的展览", size=16, color=MUTED)
    _text(s, Inches(7.1), Inches(5.2), Inches(5.5), Inches(0.35),
          "上海图书馆开放数据竞赛 · 产品介绍", size=13, color=OCHRE)
    _text(s, Inches(0.7), H - Inches(0.85), Inches(12), Inches(0.35),
          "总分总叙事  ·  概念 → 产品与技术亮点 → 内容产出", size=13, color=RICE)
    pnum()

    # ── 2 Agenda ──
    s = new_slide(prs)
    _section_title(s, "目录 · AGENDA", "总分总，三幕式讲清红鸢", "先立概念，再展开亮点，最后落到可带走的成书产出")
    agenda = [
        ("01", "总 · 概念", "问题、定位、差异化一句话", "ppt-icon-kite.png"),
        ("02", "分 · 产品亮点", "出题到成书的完整体验", "ppt-icon-walk.png"),
        ("03", "分 · 技术亮点", "防幻觉工程与策展管线", "ppt-icon-pipeline.png"),
        ("04", "总 · 内容产出", "书页 / 地图 / PDF·EPUB / 小程序", "ppt-icon-book.png"),
    ]
    for i, (num, title, desc, icon) in enumerate(agenda):
        x = Inches(0.6 + i * 3.1)
        _card(s, x, Inches(2.1), Inches(2.9), Inches(4.0))
        _icon(s, icon, x + Inches(0.95), Inches(2.45), Inches(0.9))
        _text(s, x + Inches(0.2), Inches(3.55), Inches(2.5), Inches(0.35), num, size=14, color=VERMILION, bold=True, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.15), Inches(3.95), Inches(2.6), Inches(0.5), title, size=18, bold=True, color=INK_DEEP, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.2), Inches(4.55), Inches(2.5), Inches(1.0), desc, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s, pnum())

    # ── 3 Concept one-liner ──
    s = new_slide(prs)
    _section_title(s, "总 · 概念", "别人拼接点位，我们编织关系")
    _card(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(2.0))
    _text(s, Inches(0.95), Inches(2.15), Inches(11.4), Inches(1.5),
          "用上海图书馆的真实馆藏，为你策一场只属于你的城市展览——\n走得动、站站可溯源、装订成一本书。",
          size=24, bold=True, color=INK_DEEP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    pills = [
        ("馆藏取证", "ppt-icon-library.png"),
        ("关系编织", "ppt-icon-pipeline.png"),
        ("可走展览", "ppt-icon-map.png"),
        ("成书带走", "ppt-icon-book.png"),
    ]
    for i, (label, icon) in enumerate(pills):
        x = Inches(1.1 + i * 3.0)
        _card(s, x, Inches(4.4), Inches(2.6), Inches(1.7))
        _icon(s, icon, x + Inches(0.9), Inches(4.55), Inches(0.65))
        _text(s, x + Inches(0.15), Inches(5.4), Inches(2.3), Inches(0.4), label, size=16, bold=True, color=VERMILION, align=PP_ALIGN.CENTER)
    _footer(s, pnum(), "总 · 概念")

    # ── 4 Problem ──
    s = new_slide(prs)
    _section_title(s, "总 · 概念", "红色研学不缺内容，缺的是策展力")
    problems = [
        ("点位串联", "地图上把点连起来，彼此孤立；拿到的是清单，不是故事。", "ppt-icon-map.png"),
        ("讲解难规模化", "资深策展人能讲「一栋楼的多重人生」，但人力无法批量复制。", "ppt-icon-walk.png"),
        ("AI 不可信", "通用大模型张口就来年份与人物，结论没有出处，错了你背锅。", "ppt-icon-trust.png"),
    ]
    for i, (t, d, icon) in enumerate(problems):
        y = Inches(1.9 + i * 1.55)
        _card(s, Inches(0.6), y, Inches(12.1), Inches(1.4))
        _icon(s, icon, Inches(0.9), y + Inches(0.3), Inches(0.8))
        _text(s, Inches(2.0), y + Inches(0.25), Inches(10), Inches(0.4), t, size=20, bold=True, color=VERMILION)
        _text(s, Inches(2.0), y + Inches(0.7), Inches(10.2), Inches(0.5), d, size=15, color=MUTED)
    _footer(s, pnum(), "总 · 概念")

    # ── 5 Positioning ──
    s = new_slide(prs)
    _section_title(s, "总 · 概念", "定位：AI 城市文化漫步策展", "一句话触发 · 证据先于叙事 · 成书后再翻开")
    rows = [
        ("维度", "普通路线 App / 通用大模型", "红鸢 RedTrip"),
        ("内容组织", "点位串联，彼此孤立", "建筑 × 人物 × 事件，编织关系"),
        ("可信度", "可能编造、无出处", "证据分级 A–E，找不到就明说"),
        ("可解释性", "黑盒生成", "句末上标即证据 + 反方复核"),
        ("成果形态", "页面即看即忘", "书页 / 地图 / PDF / EPUB / 小程序"),
    ]
    y0 = Inches(1.85)
    for i, row in enumerate(rows):
        y = y0 + Inches(i * 0.85)
        bg = VERMILION if i == 0 else (WHITE if i % 2 else XUAN)
        c1, c2, c3 = (PAPER, PAPER, PAPER) if i == 0 else (INK_DEEP, MUTED, INK_DEEP)
        _rect(s, Inches(0.6), y, Inches(2.4), Inches(0.75), bg if i == 0 else XUAN)
        _rect(s, Inches(3.0), y, Inches(4.6), Inches(0.75), bg)
        _rect(s, Inches(7.6), y, Inches(5.1), Inches(0.75), bg)
        _text(s, Inches(0.75), y + Inches(0.18), Inches(2.1), Inches(0.45), row[0], size=14, bold=True, color=c1 if i == 0 else VERMILION)
        _text(s, Inches(3.15), y + Inches(0.18), Inches(4.3), Inches(0.45), row[1], size=14, color=c2)
        _text(s, Inches(7.75), y + Inches(0.18), Inches(4.8), Inches(0.45), row[2], size=14, bold=(i > 0), color=c3)
    _footer(s, pnum(), "总 · 概念")

    # ── 6 Journey ──
    s = new_slide(prs)
    _section_title(s, "分 · 产品亮点", "一次策展的七步旅程")
    steps = [
        ("出题", "ppt-icon-walk.png", "城市·起点·时长·调性"),
        ("策展", "ppt-icon-pipeline.png", "取证·抽签·装订"),
        ("序章", "ppt-icon-seal.png", "主题·人物·脉络"),
        ("阅读", "ppt-icon-book.png", "站站流动散文"),
        ("反方", "ppt-icon-review.png", "策展留白挑刺"),
        ("舆图", "ppt-icon-map.png", "2.5D 可点可读"),
        ("带走", "ppt-icon-export.png", "PDF / EPUB"),
    ]
    for i, (title, icon, desc) in enumerate(steps):
        x = Inches(0.35 + i * 1.85)
        _card(s, x, Inches(2.2), Inches(1.7), Inches(3.8))
        _text(s, x + Inches(0.1), Inches(2.4), Inches(1.5), Inches(0.3), f"0{i+1}", size=12, color=OCHRE, bold=True, align=PP_ALIGN.CENTER)
        _icon(s, icon, x + Inches(0.45), Inches(2.85), Inches(0.8))
        _text(s, x + Inches(0.1), Inches(3.9), Inches(1.5), Inches(0.4), title, size=18, bold=True, color=VERMILION, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.1), Inches(4.45), Inches(1.5), Inches(1.0), desc, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s, pnum(), "分 · 产品")

    # ── 7 Product highlights grid ──
    s = new_slide(prs)
    _section_title(s, "分 · 产品亮点", "六大产品亮点一览")
    items = [
        ("证据分级 A–E", "ppt-icon-evidence.png", "先分级再写作，料薄处收着说，绝不硬编。"),
        ("可走的长散文", "ppt-icon-book.png", "每站流动叙事：进入·叠层·人物·转场。"),
        ("反方策展人", "ppt-icon-review.png", "独立挑刺：错在哪、谁缺席、哪站被跳过。"),
        ("书籍化成册", "ppt-icon-export.png", "同一叙事导出 PDF / EPUB，收藏带走。"),
        ("成书后再翻开", "ppt-icon-seal.png", "进度如实汇报，避免半成品抢先入眼。"),
        ("红鸢抽签读法", "ppt-icon-kite.png", "每次语气不同，但绝不新增史实。"),
    ]
    for i, (title, icon, desc) in enumerate(items):
        col, row = i % 3, i // 3
        x, y = Inches(0.55 + col * 4.2), Inches(1.9 + row * 2.4)
        _card(s, x, y, Inches(3.95), Inches(2.15))
        _icon(s, icon, x + Inches(0.25), y + Inches(0.35), Inches(0.7))
        _text(s, x + Inches(1.1), y + Inches(0.4), Inches(2.6), Inches(0.45), title, size=18, bold=True, color=INK_DEEP)
        _text(s, x + Inches(0.3), y + Inches(1.2), Inches(3.35), Inches(0.7), desc, size=13, color=MUTED)
    _footer(s, pnum(), "分 · 产品")

    # ── 8 Product detail: evidence ──
    s = new_slide(prs)
    _section_title(s, "分 · 产品亮点", "证据分级：把「防幻觉」做成产品体验")
    grades = [
        ("A", "一手馆藏档案", VERMILION),
        ("B", "可靠策展词库", OCHRE),
        ("C", "通用二手资料", SLATE),
        ("D/E", "传闻与推测", MUTED),
    ]
    for i, (g, label, color) in enumerate(grades):
        x = Inches(0.6 + i * 3.15)
        _card(s, x, Inches(2.0), Inches(3.0), Inches(2.4))
        seal = _round(s, x + Inches(1.0), Inches(2.25), Inches(1.0), Inches(1.0), color)
        _text(s, x + Inches(1.0), Inches(2.45), Inches(1.0), Inches(0.6), g, size=22, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.2), Inches(3.5), Inches(2.6), Inches(0.6), label, size=15, bold=True, color=INK_DEEP, align=PP_ALIGN.CENTER)
    _card(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.6))
    _text(s, Inches(0.95), Inches(4.95), Inches(11.4), Inches(1.1),
          "句末角标可点开出处 · 找不到依据就如实说「暂无数据支撑」\n数字不可被幻觉辜负 —— AI 该算清楚，而不是侃侃而谈。",
          size=16, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s, pnum(), "分 · 产品")

    # ── 9 Product: channels ──
    s = new_slide(prs)
    _section_title(s, "分 · 产品亮点", "多端触达：Web + 微信小程序", "同一策展 API，访客无需登录即可完整体验")
    _card(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.3))
    _icon(s, "ppt-icon-book.png", Inches(1.0), Inches(2.35), Inches(0.85))
    _text(s, Inches(2.0), Inches(2.45), Inches(4.2), Inches(0.5), "Web 完整书页", size=22, bold=True, color=VERMILION)
    for j, line in enumerate(["出题 Brief · SSE 真进度", "序章 / 章节阅读 / 反方留白", "2.5D 舆图 · PDF / EPUB 导出", "登录可选 · BYOK 模型配置"]):
        _text(s, Inches(1.1), Inches(3.3 + 0.55 * j), Inches(5), Inches(0.45), "·  " + line, size=15, color=INK_DEEP)

    _card(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(4.3))
    _icon(s, "ppt-icon-miniapp.png", Inches(7.2), Inches(2.35), Inches(0.85))
    _text(s, Inches(8.2), Inches(2.45), Inches(4.2), Inches(0.5), "微信小程序 MVP", size=22, bold=True, color=VERMILION)
    for j, line in enumerate(["AppID 已上架开发版 0.1.1", "出题 / 装订进度 / 序章 / 阅读", "轮询策展（适配无 EventSource）", "大赛 Token 服务端默认 GLM-5.2"]):
        _text(s, Inches(7.3), Inches(3.3 + 0.55 * j), Inches(5), Inches(0.45), "·  " + line, size=15, color=INK_DEEP)
    _footer(s, pnum(), "分 · 产品")

    # ── 10 Tech overview ──
    s = new_slide(prs)
    _section_title(s, "分 · 技术亮点", "内容管线：唯一真身", "intent → evidence → plan → narrative → polish → gate → review")
    stages = [
        ("Intent", "出题解析"),
        ("Evidence", "证据分级"),
        ("Plan", "路线规划"),
        ("Narrative", "叙事编织"),
        ("Polish", "长文润色"),
        ("Gate", "质检门禁"),
        ("Review", "反方复核"),
    ]
    for i, (en, zh) in enumerate(stages):
        x = Inches(0.4 + i * 1.85)
        _card(s, x, Inches(2.3), Inches(1.7), Inches(2.2))
        _text(s, x + Inches(0.1), Inches(2.6), Inches(1.5), Inches(0.4), f"{i+1:02d}", size=12, color=OCHRE, bold=True, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.05), Inches(3.15), Inches(1.6), Inches(0.45), en, size=14, bold=True, color=VERMILION, align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.05), Inches(3.7), Inches(1.6), Inches(0.45), zh, size=14, color=INK_DEEP, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            _text(s, x + Inches(1.55), Inches(3.15), Inches(0.35), Inches(0.4), "→", size=16, color=SLATE, align=PP_ALIGN.CENTER)
    _card(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.5))
    _text(s, Inches(0.95), Inches(5.2), Inches(11.4), Inches(1.0),
          "年份防编造是硬规则 · Gate 拦住套话与空转场 · Review 非阻断但必须可见\n降级诚实：无 key / 未过门 → status=degraded，不把假 demo 硬塞成成功。",
          size=15, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s, pnum(), "分 · 技术")

    # ── 11 Tech highlights ──
    s = new_slide(prs)
    _section_title(s, "分 · 技术亮点", "工程护城河")
    techs = [
        ("证据工程", "ppt-icon-evidence.png", ["A–E 分级写入契约", "句级 sentence_provenance", "buri / 上图 URI 可点核"]),
        ("Gate 红队", "ppt-icon-trust.png", ["字段·禁词·避坑门禁", "CI 跑 redteam runner", "失败走 degraded 不装成功"]),
        ("流式装订", "ppt-icon-pipeline.png", ["POST /curate/start", "SSE 真进度 + 小程序轮询", "成书后再呈现书页"]),
        ("多源取证", "ppt-icon-library.png", ["上海图书馆开放数据", "白名单 / 热词 / 走廊", "OSM 足迹服务地图"]),
    ]
    for i, (title, icon, bullets) in enumerate(techs):
        x = Inches(0.45 + i * 3.2)
        _card(s, x, Inches(1.95), Inches(3.05), Inches(4.5))
        _icon(s, icon, x + Inches(1.1), Inches(2.2), Inches(0.75))
        _text(s, x + Inches(0.2), Inches(3.15), Inches(2.65), Inches(0.45), title, size=18, bold=True, color=VERMILION, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            _text(s, x + Inches(0.3), Inches(3.8 + j * 0.55), Inches(2.5), Inches(0.5), "· " + b, size=13, color=INK_DEEP)
    _footer(s, pnum(), "分 · 技术")

    # ── 12 Stack ──
    s = new_slide(prs)
    _section_title(s, "分 · 技术亮点", "技术栈与仓库结构")
    left = [
        ("前端", "React 19 · TypeScript · Vite\nZustand · XState · Three.js 舆图"),
        ("后端", "FastAPI · Python 策展管线\nJWT 鉴权 · BYOK 加密存储"),
        ("交付", "Web /redtrip · 微信小程序\nsy-realm.ltd 线上可演示"),
    ]
    for i, (h, body) in enumerate(left):
        y = Inches(1.95 + i * 1.55)
        _card(s, Inches(0.6), y, Inches(6.0), Inches(1.4))
        _text(s, Inches(0.9), y + Inches(0.2), Inches(5.4), Inches(0.35), h, size=16, bold=True, color=VERMILION)
        _text(s, Inches(0.9), y + Inches(0.6), Inches(5.4), Inches(0.65), body, size=14, color=INK_DEEP)

    _card(s, Inches(6.9), Inches(1.95), Inches(5.8), Inches(4.55))
    _text(s, Inches(7.2), Inches(2.2), Inches(5.2), Inches(0.4), "Monorepo", size=16, bold=True, color=VERMILION)
    mono = (
        "apps/web          书页前端\n"
        "apps/api          FastAPI\n"
        "apps/miniprogram  微信小程序\n"
        "packages/curator  策展管线\n"
        "packages/gate     质检 / 红队\n"
        "packages/library-client  上图等客户端\n"
        "packages/contracts 共享契约\n"
        "content/          白名单·fixtures·热词"
    )
    _text(s, Inches(7.2), Inches(2.75), Inches(5.2), Inches(3.4), mono, size=14, color=INK_DEEP)
    _footer(s, pnum(), "分 · 技术")

    # ── 13 Output overview (总) ──
    s = new_slide(prs)
    _section_title(s, "总 · 内容产出", "策展的终点，是一本可以带走的书")
    outs = [
        ("序章书页", "ppt-icon-seal.png", "主题 · 人物 · 章节脉络 · 红鸢读法"),
        ("站站章节", "ppt-icon-book.png", "流动散文 · 句级溯源 · 行前提示"),
        ("叙事舆图", "ppt-icon-map.png", "2.5D 海派配色 · 点站即读"),
        ("导出成册", "ppt-icon-export.png", "PDF / EPUB · 收藏与分享"),
        ("反方附录", "ppt-icon-review.png", "策展留白 · 敢于摊开犹豫"),
        ("小程序书", "ppt-icon-miniapp.png", "手机端完整体验切片"),
    ]
    for i, (title, icon, desc) in enumerate(outs):
        col, row = i % 3, i // 3
        x, y = Inches(0.55 + col * 4.2), Inches(1.95 + row * 2.35)
        _card(s, x, y, Inches(3.95), Inches(2.15))
        _icon(s, icon, x + Inches(0.25), y + Inches(0.4), Inches(0.7))
        _text(s, x + Inches(1.15), y + Inches(0.45), Inches(2.55), Inches(0.45), title, size=18, bold=True, color=INK_DEEP)
        _text(s, x + Inches(0.3), y + Inches(1.25), Inches(3.35), Inches(0.6), desc, size=13, color=MUTED)
    _footer(s, pnum(), "总 · 产出")

    # ── 14 Demo content ──
    s = new_slide(prs)
    _section_title(s, "总 · 内容产出", "竞赛演示线：秒开看成书形态")
    _card(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(4.4))
    _icon(s, "ppt-icon-walk.png", Inches(1.0), Inches(2.3), Inches(0.75))
    _text(s, Inches(1.95), Inches(2.4), Inches(4.3), Inches(0.45), "演示武康 · 六站可溯源", size=20, bold=True, color=VERMILION)
    _text(s, Inches(1.0), Inches(3.3), Inches(5.2), Inches(2.6),
          "冻结包，不等待 LLM\n巴金故居 → 周璇旧居 → 武康大楼\n→ 宋庆龄故居 → 丁香花园 → 武康庭\n\n每站可点上图 buri · 句级出处可演示",
          size=15, color=INK_DEEP)

    _card(s, Inches(6.9), Inches(1.95), Inches(5.8), Inches(4.4))
    _icon(s, "ppt-icon-map.png", Inches(7.3), Inches(2.3), Inches(0.75))
    _text(s, Inches(8.25), Inches(2.4), Inches(4.1), Inches(0.45), "演示一大·外滩", size=20, bold=True, color=VERMILION)
    _text(s, Inches(7.3), Inches(3.3), Inches(5.0), Inches(2.6),
          "通道诚实叙事样板\n展示「馆藏 / 词库 / OSM」分层\n\n适合向评委说明：\n我们如何标注证据通道，而不是假装全能。",
          size=15, color=INK_DEEP)
    _footer(s, pnum(), "总 · 产出")

    # ── 15 Closing ──
    s = new_slide(prs)
    _rect(s, 0, 0, W, H, INK_DEEP)
    _rect(s, 0, 0, W, Inches(0.12), VERMILION)
    _icon(s, "ppt-icon-kite.png", Inches(6.05), Inches(1.3), Inches(1.2))
    _text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.7), "数字不可被幻觉辜负", size=36, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
          "AI 该算清楚，而不是侃侃而谈。\n红鸢把馆藏变成一场走得动的展览。",
          size=18, color=RICE, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.9),
          "Web  https://sy-realm.ltd/redtrip/\nGitHub  github.com/MiLab-Bit/RedTrip\n小程序  AppID wxc7953007477c1980",
          size=14, color=OCHRE, align=PP_ALIGN.CENTER)
    _text(s, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.4),
          "谢谢 · 欢迎提问", size=16, color=SLATE, align=PP_ALIGN.CENTER)
    pnum()

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    ART.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ART))
    print(f"wrote {OUT}")
    print(f"wrote {ART}")
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
